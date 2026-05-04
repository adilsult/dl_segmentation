"""Build BTS_Final_Presentation.pptx from the AI700-001 official template.

Strategy:
- Open AI700-001_BK_SPRG2026-PPT-Template-Final.pptx (preserves design, headers,
  background, slide numbers).
- For each section, edit the matching template slide. Where one section needs
  more space than one slide, duplicate the slide.
- Render math formulas as PNGs via matplotlib mathtext (LaTeX look) and embed.
- Image slots remain as labeled grey placeholder rectangles for the user.

Run: python3 build_presentation.py
"""

from copy import deepcopy
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).parent
TEMPLATE = ROOT / "AI700-001_BK_SPRG2026-PPT-Template-Final.pptx"
OUT_PATH = ROOT / "BTS_Final_Presentation.pptx"
FORMULA_DIR = ROOT / "_formulas"
FORMULA_DIR.mkdir(exist_ok=True)

NEAR_BLACK = RGBColor(0x21, 0x21, 0x21)
NAVY = RGBColor(0x0F, 0x2A, 0x6E)
ACCENT = RGBColor(0xC8, 0x16, 0x2C)  # template's red accent (visually similar)
GREY_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREY_BORDER = RGBColor(0xBD, 0xBD, 0xBD)
GREY_CAPTION = RGBColor(0x61, 0x61, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xFF, 0xF4, 0xCE)

BODY_FONT = "Calibri"
TITLE_FONT = "Calibri"
MONO_FONT = "Consolas"


# ---------- Formula rendering ----------
plt.rcParams["mathtext.fontset"] = "cm"  # Computer Modern (LaTeX look)
plt.rcParams["mathtext.rm"] = "serif"


def render_formula(latex, filename, fontsize=22):
    """Render a LaTeX math string to a transparent PNG."""
    path = FORMULA_DIR / filename
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(0, 0, f"${latex}$", fontsize=fontsize, color="black")
    fig.savefig(path, dpi=300, bbox_inches="tight", pad_inches=0.05,
                transparent=True)
    plt.close(fig)
    return path


# Pre-render all formulas
F = {}
F["dicebce"] = render_formula(
    r"\mathcal{L}_{DiceBCE} = \frac{1}{2}\,\mathcal{L}_{BCE} + \frac{1}{2}\,\mathcal{L}_{Dice}",
    "dicebce.png", fontsize=24)
F["bce"] = render_formula(
    r"\mathcal{L}_{BCE} = -\frac{1}{N}\sum_{i}\left[y_i \log p_i + (1-y_i)\log(1-p_i)\right]",
    "bce.png", fontsize=22)
F["dice"] = render_formula(
    r"\mathcal{L}_{Dice} = 1 - \frac{2\sum_i p_i y_i + \epsilon}{\sum_i p_i + \sum_i y_i + \epsilon}",
    "dice.png", fontsize=24)
F["tversky"] = render_formula(
    r"TI = \frac{TP + \epsilon}{TP + \alpha\,FP + \beta\,FN + \epsilon}",
    "tversky.png", fontsize=24)
F["focaltversky"] = render_formula(
    r"\mathcal{L}_{FT} = (1 - TI)^{\gamma}",
    "focaltversky.png", fontsize=26)
F["attention"] = render_formula(
    r"\alpha = \sigma\left(W_{\psi}\cdot \mathrm{ReLU}\left(W_g(g) + W_x(x)\right)\right),"
    r"\quad \tilde{x} = x \odot \alpha",
    "attention.png", fontsize=22)
F["adam_m"] = render_formula(
    r"m_t = \beta_1 m_{t-1} + (1-\beta_1)\,g_t,\quad v_t = \beta_2 v_{t-1} + (1-\beta_2)\,g_t^{2}",
    "adam_m.png", fontsize=22)
F["adam_hat"] = render_formula(
    r"\hat{m}_t = \frac{m_t}{1-\beta_1^{t}},\quad \hat{v}_t = \frac{v_t}{1-\beta_2^{t}}",
    "adam_hat.png", fontsize=22)
F["adam_update"] = render_formula(
    r"\theta_t = \theta_{t-1} - \eta\,\frac{\hat{m}_t}{\sqrt{\hat{v}_t}+\epsilon}",
    "adam_update.png", fontsize=24)
F["dice_metric"] = render_formula(
    r"\mathrm{Dice} = \frac{2\,TP + \epsilon}{2\,TP + FP + FN + \epsilon}",
    "dice_metric.png", fontsize=24)
F["iou_metric"] = render_formula(
    r"\mathrm{IoU} = \frac{TP + \epsilon}{TP + FP + FN + \epsilon}",
    "iou_metric.png", fontsize=24)
F["dice_iou_eq"] = render_formula(
    r"\mathrm{Dice} = \frac{2\,\mathrm{IoU}}{1 + \mathrm{IoU}}",
    "dice_iou_eq.png", fontsize=24)


# ---------- Slide manipulation helpers ----------
def duplicate_slide(prs, src_index):
    """Append a deep copy of slides[src_index] to the end of the deck."""
    src = prs.slides[src_index]
    new = prs.slides.add_slide(src.slide_layout)
    # remove default placeholders that came with the layout
    for shp in list(new.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    # copy all shapes from source
    for shp in src.shapes:
        new.shapes._spTree.append(deepcopy(shp._element))
    return new


def move_slide(prs, old_idx, new_idx):
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    item = items[old_idx]
    sldIdLst.remove(item)
    sldIdLst.insert(new_idx, item)


def find_shape(slide, name_substr):
    for shp in slide.shapes:
        if name_substr.lower() in shp.name.lower():
            return shp
    return None


def find_content_placeholder(slide):
    """Return the Content/Text Placeholder shape on a slide (the big editable area)."""
    candidates = []
    for shp in slide.shapes:
        if "content placeholder" in shp.name.lower() or \
           "text placeholder" in shp.name.lower():
            candidates.append(shp)
    if not candidates:
        return None
    # prefer the largest
    return max(candidates, key=lambda s: (s.width or 0) * (s.height or 0))


def clear_text_frame(tf):
    """Remove all paragraphs except a single empty one."""
    txBody = tf._txBody
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    # add empty p back
    p = etree.SubElement(txBody, qn("a:p"))


def set_paragraphs(tf, paragraphs, *, font=BODY_FONT, default_size=16,
                   color=NEAR_BLACK):
    """paragraphs = list; each item is either a string (uses defaults) or
    a dict {text, size, bold, italic, color, font, level, bullet}.
    """
    clear_text_frame(tf)
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        if isinstance(item, str):
            item = {"text": item}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = item.get("align", PP_ALIGN.LEFT)
        if "level" in item:
            p.level = item["level"]
        run = p.add_run()
        run.text = item.get("text", "")
        run.font.name = item.get("font", font)
        run.font.size = Pt(item.get("size", default_size))
        run.font.bold = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.color.rgb = item.get("color", color)


def add_textbox(slide, left, top, width, height, paragraphs, **kw):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if isinstance(paragraphs, str):
        paragraphs = [paragraphs]
    set_paragraphs(tf, paragraphs, **kw)
    return tb


def add_placeholder_box(slide, left, top, width, height, caption):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREY_BG
    shape.line.color.rgb = GREY_BORDER
    shape.line.width = Pt(1.25)
    ln = shape.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    paragraphs = [{"text": ln, "size": 11, "italic": True,
                   "color": GREY_CAPTION, "align": PP_ALIGN.CENTER}
                  for ln in caption.split("\n")]
    set_paragraphs(tf, paragraphs)
    return shape


def add_table(slide, left, top, width, height, data, *,
              header_fill=NAVY, header_color=WHITE, body_size=11,
              header_size=11, col_widths=None, highlight_row=None):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.name = BODY_FONT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.color.rgb = header_color
                run.font.bold = True
                run.font.size = Pt(header_size)
            else:
                run.font.size = Pt(body_size)
                run.font.color.rgb = NEAR_BLACK
                cell.fill.solid()
                if highlight_row is not None and r == highlight_row:
                    cell.fill.fore_color.rgb = HIGHLIGHT
                    run.font.bold = True
                else:
                    cell.fill.fore_color.rgb = WHITE
    return tbl


def add_formula(slide, key, left, top, height):
    """Embed a pre-rendered formula PNG, anchoring by height."""
    return slide.shapes.add_picture(str(F[key]), left, top, height=height)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


# ---------- Build ----------
prs = Presentation(str(TEMPLATE))

# ===== Slide 1: Title =====
s = prs.slides[0]
# Rectangle 2 = project title (top-left red bar area)
title_rect = find_shape(s, "Rectangle 2")
set_paragraphs(title_rect.text_frame, [
    {"text": "Brain Tumor Segmentation", "size": 28, "bold": True, "color": WHITE},
    {"text": "from Multi-Modal MRI", "size": 24, "bold": True, "color": WHITE},
], font=TITLE_FONT)
# Rectangle 3 = subtitle area (lower-left big block)
sub_rect = find_shape(s, "Rectangle 3")
if sub_rect is not None:
    set_paragraphs(sub_rect.text_frame, [
        {"text": "A Comparative Study of 2D U-Net Variants",
         "size": 18, "bold": True, "color": WHITE},
        {"text": "with Attention Gates and Transformer Bottleneck",
         "size": 16, "italic": True, "color": WHITE},
    ], font=TITLE_FONT)
# TextBox 8 = By: name | date | course
authors_tb = find_shape(s, "TextBox 8")
if authors_tb is not None:
    set_paragraphs(authors_tb.text_frame, [
        {"text": "By: Adilsultan Khairolla, Kavyashree M. Vijaykumar, Reda",
         "size": 11, "color": NEAR_BLACK},
        {"text": "Date: May 2026", "size": 11, "color": NEAR_BLACK},
        {"text": "AI700-001_BK_Spring 2026", "size": 11, "color": NEAR_BLACK},
    ])
# TextBox 5 = "Picture Representing The Scope Of Your Project"
pic_tb = find_shape(s, "TextBox 5")
if pic_tb is not None:
    set_paragraphs(pic_tb.text_frame, [
        {"text": "[Insert: project hero image]", "size": 14,
         "italic": True, "color": GREY_CAPTION, "align": PP_ALIGN.CENTER},
        {"text": "(e.g., T1ce panel from BraTS preview)", "size": 11,
         "italic": True, "color": GREY_CAPTION, "align": PP_ALIGN.CENTER},
    ])


# ===== Slide 2: Outlines =====
s = prs.slides[1]
cp = find_content_placeholder(s)
outline_items = [
    "1.  Introduction & Clinical Motivation",
    "2.  Application Domain",
    "3.  Literature Review — Foundations + Recent SOTA",
    "4.  Comparison with Prior Work",
    "5.  Methodology — Dataset, Architectures, Loss, Optimizer",
    "6.  Mathematical Framework — Loss & Metric Formulas",
    "7.  Results — Quantitative & Qualitative",
    "8.  Discussion — Capacity vs. Data Tradeoff",
    "9.  Conclusion & Future Work",
    "10. References",
]
set_paragraphs(cp.text_frame,
               [{"text": t, "size": 20} for t in outline_items])


# ===== Slide 3: Introduction =====
s = prs.slides[2]
cp = find_content_placeholder(s)
intro = [
    {"text": "Why automate brain tumor segmentation?", "size": 18, "bold": True, "color": NAVY},
    {"text": "", "size": 6},
    {"text": "• Gliomas: irregular shapes, ambiguous boundaries, high inter-patient variability", "size": 14},
    {"text": "• Manual segmentation is slow (~hours/scan); not scalable to population screening", "size": 14},
    {"text": "• Inter-rater Dice agreement between expert radiologists: 0.88–0.92", "size": 14},
    {"text": "• Automated methods → reproducibility and 24/7 availability", "size": 14},
    {"text": "• BraTS 2020 = standardized benchmark: 369 patients, 4 MRI modalities (T1 / T1ce / T2 / FLAIR)", "size": 14},
    {"text": "", "size": 6},
    {"text": "Project goal", "size": 16, "bold": True, "color": ACCENT},
    {"text": "Quantify how each architectural component — skip connections, spatial attention, "
             "and global self-attention — contributes to segmentation accuracy.",
     "size": 14, "italic": True},
]
set_paragraphs(cp.text_frame, intro)


# ===== Slide 4: Application =====
s = prs.slides[3]
cp = find_content_placeholder(s)
app_items = [
    {"text": "Where does this technology matter?", "size": 18, "bold": True, "color": NAVY},
    {"text": "", "size": 6},
    {"text": "Diagnosis", "size": 15, "bold": True, "color": ACCENT},
    {"text": "    Quantitative tumor volume → grading (low- vs. high-grade glioma)", "size": 13},
    {"text": "Treatment Planning", "size": 15, "bold": True, "color": ACCENT},
    {"text": "    Precise mask drives radiation oncology dose targeting (proton/photon)", "size": 13},
    {"text": "Treatment Monitoring", "size": 15, "bold": True, "color": ACCENT},
    {"text": "    Track tumor growth/shrinkage between scans (RANO criteria)", "size": 13},
    {"text": "Surgical Navigation", "size": 15, "bold": True, "color": ACCENT},
    {"text": "    Integrate with neuronavigation systems for resection guidance", "size": 13},
    {"text": "Research Acceleration", "size": 15, "bold": True, "color": ACCENT},
    {"text": "    Enables population-scale neuro-oncology studies", "size": 13},
    {"text": "", "size": 6},
    {"text": "Clinical adoption requires Dice ≥ 0.85 and low Hausdorff distance.",
     "size": 13, "italic": True, "color": NAVY},
]
set_paragraphs(cp.text_frame, app_items)


# ===== Slide 5: Literature Review =====
s = prs.slides[4]
cp = find_content_placeholder(s)
lit = [
    {"text": "From U-Net to Transformers — A Decade of Medical Segmentation",
     "size": 17, "bold": True, "color": NAVY},
    {"text": "", "size": 6},
    {"text": "Architectural lineage", "size": 14, "bold": True, "color": ACCENT},
    {"text": "• U-Net (Ronneberger et al., MICCAI 2015) — symmetric encoder-decoder + skip connections; foundational", "size": 12},
    {"text": "• V-Net (Milletari et al., 3DV 2016) — 3D extension; introduced Dice Loss for class imbalance", "size": 12},
    {"text": "• Attention U-Net (Oktay et al., MIDL 2018) — soft attention gates suppress irrelevant skip features", "size": 12},
    {"text": "• TransUNet (Chen et al., 2021) — first hybrid CNN–Transformer for medical images", "size": 12},
    {"text": "• TransBTS (Wang et al., MICCAI 2021) — Transformer at bottleneck for brain tumor segmentation", "size": 12},
    {"text": "• Swin UNETR (Hatamizadeh et al., 2022) — hierarchical shifted-window attention", "size": 12},
    {"text": "• nnU-Net (Isensee et al., Nat. Methods 2021) — self-configuring framework, BraTS winner", "size": 12},
    {"text": "", "size": 6},
    {"text": "Loss function evolution", "size": 14, "bold": True, "color": ACCENT},
    {"text": "BCE  →  Dice (Milletari, 2016)  →  Tversky (Salehi, 2017)  →  Focal Tversky (Abraham & Khan, 2019)", "size": 12},
    {"text": "Each step targets class imbalance more aggressively.", "size": 11, "italic": True},
    {"text": "", "size": 4},
    {"text": "Our Hybrid model adapts TransBTS's bottleneck idea to 2D for computational efficiency.",
     "size": 11, "italic": True, "color": NAVY},
]
set_paragraphs(cp.text_frame, lit)


# ===== Slide 6: Literature Review Results =====
s = prs.slides[5]
# replace the small content placeholder with table + bullets
cp = find_content_placeholder(s)
if cp is not None:
    remove_shape(cp)
# Table
table_data = [
    ["Method", "Dim", "Params", "Dice (WT)", "Year"],
    ["U-Net (Ronneberger)", "2D", "~31M", "0.85–0.88", "2015"],
    ["Attention U-Net (Oktay)", "2D/3D", "~2–30M", "0.87–0.89", "2018"],
    ["TransBTS (Wang)", "3D", "~30M", "0.89–0.90", "2021"],
    ["Swin UNETR (Hatamizadeh)", "3D", "~62M", "0.90–0.91", "2022"],
    ["nnU-Net (Isensee)", "3D", "~30M", "0.91–0.92", "2021+"],
    ["Our 2D Hybrid (lightweight)", "2D", "2.84M", "0.83 (50 pts)", "2026"],
]
add_table(s, Inches(0.5), Inches(1.15), Inches(9.0), Inches(2.7), table_data,
          col_widths=[Inches(2.9), Inches(1.0), Inches(1.5), Inches(2.2), Inches(1.4)],
          highlight_row=6, body_size=11, header_size=11)
# Key bullets
add_textbox(s, Inches(0.5), Inches(4.05), Inches(9.0), Inches(2.7), [
    {"text": "Key positioning", "size": 14, "bold": True, "color": ACCENT},
    {"text": "• Our 2D approach is 10–20× lighter than 3D SOTA", "size": 12},
    {"text": "• Trains in ~15 min on T4 GPU vs. days on A100 for nnU-Net", "size": 12},
    {"text": "• Result is competitive with typical 2D U-Net papers (0.80–0.85)", "size": 12},
    {"text": "• 2D-vs-3D gap (~0.05–0.07 Dice) dominates over architecture choice within 2D", "size": 12},
])


# ===== Slide 7: Methodology — Dataset =====
s = prs.slides[6]
# kill default content placeholders; rebuild
for shp in list(s.shapes):
    if "placeholder" in shp.name.lower() and "slide number" not in shp.name.lower():
        remove_shape(shp)
# Section title under header
add_textbox(s, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
            [{"text": "Methodology #1: Dataset & Pipeline — BraTS 2020",
              "size": 18, "bold": True, "color": NAVY}])
# Left bullets
add_textbox(s, Inches(0.5), Inches(1.6), Inches(5.3), Inches(4.5), [
    {"text": "• 369 patients with histologically confirmed gliomas", "size": 13},
    {"text": "• 4 co-registered modalities (T1, T1ce, T2, FLAIR), 240×240×155 voxels", "size": 13},
    {"text": "• Per-pixel labels: NCR (1), edema (2), enhancing tumor (4)", "size": 13},
    {"text": "• Subset of 50 patients (compute budget) → ~3,200 tumor-bearing slices", "size": 13},
    {"text": "• Patient-level 80/20 split (40 train / 10 val) — prevents leakage", "size": 13},
    {"text": "• Tumor pixels ≈ 7% of slice → severe class imbalance", "size": 13},
    {"text": "", "size": 6},
    {"text": "Each modality reveals different tumor sub-structure",
     "size": 11, "italic": True, "color": NAVY},
    {"text": "(T1ce → core, FLAIR → edema). Multi-class merged to binary mask.",
     "size": 11, "italic": True, "color": NAVY},
])
# Right placeholder
add_placeholder_box(s, Inches(5.95), Inches(1.6), Inches(3.85), Inches(4.5),
                    "[Insert: cell_13_dataset_preview_mri_slices.png]\nfour modalities + tumor mask overlay")


# ===== Add 3 more methodology slides after slide 7 (UNet, Attention, Hybrid) =====
def make_methodology_slide(prs, src_idx, title, specs, params_line, placeholder_caption):
    """Duplicate the methodology slide and fill it in."""
    new = duplicate_slide(prs, src_idx)
    # remove all non-header shapes (keep title bar Rectangle 2 + slide number)
    for shp in list(new.shapes):
        nm = shp.name.lower()
        if "rectangle 2" in nm or "slide number" in nm:
            continue
        remove_shape(shp)
    # subtitle
    add_textbox(new, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
                [{"text": title, "size": 18, "bold": True, "color": NAVY}])
    # left placeholder for diagram
    add_placeholder_box(new, Inches(0.5), Inches(1.6), Inches(5.0), Inches(4.5),
                        placeholder_caption)
    # right specs
    para = [{"text": "Specifications", "size": 14, "bold": True, "color": ACCENT}]
    para += [{"text": "• " + s, "size": 12} for s in specs]
    para += [{"text": "", "size": 6},
             {"text": params_line, "size": 14, "bold": True, "color": NAVY}]
    add_textbox(new, Inches(5.65), Inches(1.6), Inches(4.15), Inches(4.5), para)
    return new


# UNet2D
make_methodology_slide(prs, 6,
    title="Methodology #2: UNet2D — Symmetric Encoder-Decoder",
    specs=[
        "DoubleConv: Conv3×3 → BN → ReLU → Conv3×3 → BN → ReLU",
        "4 Down blocks (MaxPool2 + DoubleConv)",
        "4 Up blocks (ConvTranspose2 + concat + DoubleConv)",
        "Output head: 1×1 Conv → logits",
        "Channels: 4 → 16 → 32 → 64 → 128 → 256",
        "Hidden conv layers: 18 — Total trainable: 23",
    ],
    params_line="Params (base=16):  1,942,721",
    placeholder_caption="[Insert: unet_arch.png]\nencoder–decoder block diagram\nwith channel/spatial labels and skip arrows")

# AttentionUNet2D — add slide; we'll also drop the formula image on it
att_slide = make_methodology_slide(prs, 6,
    title="Methodology #3: AttentionUNet2D — Spatial Recalibration of Skips",
    specs=[
        "4 AttentionGate modules (cheap 1×1 convs)",
        "Hidden layers: 23 + 12 = 35",
        "α≈1 ⇒ keep (likely tumor); α≈0 ⇒ suppress",
        "Reference: Oktay et al., MIDL 2018",
    ],
    params_line="Params:  1,986,965  (+44K vs UNet2D)",
    placeholder_caption="[Insert: attention_gate.png]\nAttention Gate dataflow\n(g, x → 1×1 convs → ReLU → σ → ⊙)")
# Attention formula image at bottom-left under the placeholder
add_formula(att_slide, "attention", Inches(0.5), Inches(6.2), Inches(0.45))

# HybridUNet2D
make_methodology_slide(prs, 6,
    title="Methodology #4: HybridUNet2D — Transformer Bottleneck",
    specs=[
        "Tokens: 15×15 = 225, dim 256",
        "Each layer: MHSA(4 heads) + FFN(256→256, GELU)",
        "2× LayerNorm + Dropout 0.1",
        "Learnable pos. emb. (1, 225, 256), N(0, 0.02²)",
        "2 transformer layers stacked",
        "Hidden layers: 35 + 8 = 43",
        "Self-attention → global view at bottleneck",
    ],
    params_line="Params:  2,836,629  (+850K)",
    placeholder_caption="[Insert: transformer_bottleneck.png]\nbottleneck (B,256,15,15) → flatten → 225 tokens\n→ +pos.emb → TransformerEncoder×2 → reshape")


# At this point: methodology dup slides are at end (indexes 12, 13, 14).
# Original order: [0]title [1]outline [2]intro [3]app [4]lit [5]litres [6]meth-data
# [7]understand-fn [8]results [9]conclusion [10]refs [11]questions
# After append: meth-unet=12, meth-att=13, meth-hyb=14
# We want them right after slide 7 (index 6). Move them to indexes 7, 8, 9.
move_slide(prs, 12, 7)   # UNet → 7
move_slide(prs, 13, 8)   # Attn → 8
move_slide(prs, 14, 9)   # Hybrid → 9
# Now layout:
# 0 title, 1 outline, 2 intro, 3 app, 4 lit, 5 litres, 6 meth-data,
# 7 meth-unet, 8 meth-att, 9 meth-hyb, 10 understand-fn, 11 results,
# 12 conclusion, 13 refs, 14 questions


# ===== Slide index 10: Understanding the Functions — Loss Functions =====
s = prs.slides[10]
for shp in list(s.shapes):
    nm = shp.name.lower()
    if "rectangle 2" in nm or "slide number" in nm:
        continue
    remove_shape(shp)
add_textbox(s, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
            [{"text": "Loss Functions for Class-Imbalanced Segmentation",
              "size": 18, "bold": True, "color": NAVY}])
# Left: DiceBCE
add_textbox(s, Inches(0.5), Inches(1.65), Inches(4.6), Inches(0.4),
            [{"text": "DiceBCE Loss", "size": 15, "bold": True, "color": ACCENT}])
add_formula(s, "dicebce", Inches(0.5), Inches(2.10), Inches(0.45))
add_formula(s, "bce",     Inches(0.5), Inches(2.75), Inches(0.45))
add_formula(s, "dice",    Inches(0.5), Inches(3.50), Inches(0.65))
add_textbox(s, Inches(0.5), Inches(4.6), Inches(4.6), Inches(2.0), [
    {"text": "• BCE: stable per-pixel gradients", "size": 12},
    {"text": "• Dice: direct overlap optimization,", "size": 12},
    {"text": "    robust to imbalance", "size": 12},
])
# Right: Focal Tversky
add_textbox(s, Inches(5.2), Inches(1.65), Inches(4.6), Inches(0.4),
            [{"text": "Focal Tversky Loss", "size": 15, "bold": True, "color": ACCENT}])
add_formula(s, "tversky",      Inches(5.2), Inches(2.10), Inches(0.55))
add_formula(s, "focaltversky", Inches(5.2), Inches(2.95), Inches(0.50))
ft = [
    ["Param", "Value", "Effect"],
    ["α", "0.7", "weight on FP"],
    ["β", "0.3", "weight on FN"],
    ["γ", "0.75", "focal exponent"],
]
add_table(s, Inches(5.2), Inches(3.7), Inches(4.6), Inches(1.6), ft,
          col_widths=[Inches(0.9), Inches(0.9), Inches(2.8)],
          body_size=11, header_size=11, header_fill=ACCENT)
add_textbox(s, Inches(5.2), Inches(5.5), Inches(4.6), Inches(1.2),
            [{"text": "β > α  ⇒  missing tumor (FN) penalised",
              "size": 11, "italic": True, "color": NAVY},
             {"text": ">2× more than false alarm (FP).",
              "size": 11, "italic": True, "color": NAVY}])


# ===== Add a duplicate Understanding slide for Optimizer & Metrics =====
opt_slide = duplicate_slide(prs, 10)
for shp in list(opt_slide.shapes):
    nm = shp.name.lower()
    if "rectangle 2" in nm or "slide number" in nm:
        continue
    remove_shape(shp)
add_textbox(opt_slide, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
            [{"text": "Adam Optimizer & Evaluation Metrics",
              "size": 18, "bold": True, "color": NAVY}])
# Left: Adam
add_textbox(opt_slide, Inches(0.5), Inches(1.65), Inches(4.6), Inches(0.4),
            [{"text": "Adam optimizer", "size": 15, "bold": True, "color": ACCENT}])
add_formula(opt_slide, "adam_m",      Inches(0.5), Inches(2.10), Inches(0.45))
add_formula(opt_slide, "adam_hat",    Inches(0.5), Inches(2.75), Inches(0.45))
add_formula(opt_slide, "adam_update", Inches(0.5), Inches(3.40), Inches(0.55))
add_textbox(opt_slide, Inches(0.5), Inches(4.3), Inches(4.6), Inches(2.5), [
    {"text": "Hyperparameters", "size": 13, "bold": True, "color": NAVY},
    {"text": "• lr = 1e-3,  β₁=0.9,  β₂=0.999,  ε=1e-8", "size": 11},
    {"text": "• weight_decay = 1e-5  (L2)", "size": 11},
    {"text": "• ReduceLROnPlateau(factor=0.5, patience=5)", "size": 11},
    {"text": "• Early stopping: patience 10", "size": 11},
])
# Right: Metrics
add_textbox(opt_slide, Inches(5.2), Inches(1.65), Inches(4.6), Inches(0.4),
            [{"text": "Evaluation metrics", "size": 15, "bold": True, "color": ACCENT}])
add_formula(opt_slide, "dice_metric",  Inches(5.2), Inches(2.10), Inches(0.55))
add_formula(opt_slide, "iou_metric",   Inches(5.2), Inches(2.95), Inches(0.55))
add_formula(opt_slide, "dice_iou_eq",  Inches(5.2), Inches(3.80), Inches(0.55))
add_textbox(opt_slide, Inches(5.2), Inches(4.5), Inches(4.6), Inches(2.5), [
    {"text": "Training config", "size": 13, "bold": True, "color": NAVY},
    {"text": "• Batch size 8, base channels 16, 30 epochs", "size": 11},
    {"text": "• NVIDIA T4 (15 GB), CUDA 12.8, PyTorch 2.10", "size": 11},
    {"text": "• cudnn.deterministic=True, seed 42", "size": 11},
])
# Move it right after slide 10
# After append, it's at end. Total slides currently = 16. Move from 15 → 11.
move_slide(prs, len(prs.slides) - 1, 11)
# Now: ... 10 understand-loss, 11 understand-opt, 12 results, 13 conclusion, 14 refs, 15 questions


# ===== Slide index 12: Obtained Results — Training Curves =====
s = prs.slides[12]
for shp in list(s.shapes):
    nm = shp.name.lower()
    if "rectangle 2" in nm or "slide number" in nm:
        continue
    remove_shape(shp)
add_textbox(s, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
            [{"text": "Training Curves & Convergence",
              "size": 18, "bold": True, "color": NAVY}])
add_placeholder_box(s, Inches(0.5), Inches(1.65), Inches(9.0), Inches(3.8),
                    "[Insert: cell_15_training_curves_all_configs.png]")
add_textbox(s, Inches(0.5), Inches(5.55), Inches(9.0), Inches(1.4), [
    {"text": "• All four configurations plateau at validation Dice 0.82–0.84 by epoch 18–26", "size": 11},
    {"text": "• HybridUNet2D + FocalTversky shows highest training loss (different scale; not comparable to DiceBCE)", "size": 11},
    {"text": "• DiceBCE configurations exhibit smooth monotonic improvement", "size": 11},
    {"text": "• ReduceLROnPlateau triggered ~epoch 12 (visible plateau before further improvement)", "size": 11},
])


# ===== Add duplicate Obtained Results slide: final numbers =====
res2 = duplicate_slide(prs, 12)
for shp in list(res2.shapes):
    nm = shp.name.lower()
    if "rectangle 2" in nm or "slide number" in nm:
        continue
    remove_shape(shp)
add_textbox(res2, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
            [{"text": "Final Validation Metrics",
              "size": 18, "bold": True, "color": NAVY}])
results_data = [
    ["Configuration", "Best Dice", "Best IoU", "Best Epoch", "Params"],
    ["UNet2D + DiceBCE  ★", "0.8370", "0.7715", "26", "1.94M"],
    ["AttentionUNet2D + DiceBCE", "0.8291", "0.7635", "18", "1.99M"],
    ["HybridUNet2D + DiceBCE", "0.8282", "0.7677", "22", "2.84M"],
    ["HybridUNet2D + FocalTversky", "0.8206", "0.7521", "25", "2.84M"],
]
add_table(res2, Inches(0.5), Inches(1.65), Inches(9.0), Inches(2.0), results_data,
          col_widths=[Inches(3.4), Inches(1.3), Inches(1.3), Inches(1.5), Inches(1.5)],
          highlight_row=1, body_size=12, header_size=12)
add_placeholder_box(res2, Inches(0.5), Inches(3.85), Inches(5.4), Inches(3.0),
                    "[Insert: cell_17_sample_predictions_6_slices.png]\n6 slices: FLAIR | GT | Pred")
add_textbox(res2, Inches(6.05), Inches(3.85), Inches(3.75), Inches(3.0), [
    {"text": "Observations", "size": 13, "bold": True, "color": ACCENT},
    {"text": "• Boundary delineation accurate on large tumors", "size": 11},
    {"text": "• Minor under-segmentation at edema-tissue transitions", "size": 11},
    {"text": "• No catastrophic false positives in healthy tissue", "size": 11},
])
move_slide(prs, len(prs.slides) - 1, 13)


# ===== Add Discussion slide =====
disc = duplicate_slide(prs, 12)
for shp in list(disc.shapes):
    nm = shp.name.lower()
    if "rectangle 2" in nm or "slide number" in nm:
        continue
    remove_shape(shp)
add_textbox(disc, Inches(0.69), Inches(1.05), Inches(8.6), Inches(0.5),
            [{"text": "Discussion: Why did the simpler UNet2D win?",
              "size": 18, "bold": True, "color": NAVY}])
boxes = [
    ("1.  Capacity-vs-Data Tradeoff", [
        "HybridUNet2D has 46% more params than UNet2D",
        "50 patients (~3,200 slices) too small to fully exploit transformer capacity",
        "\"Bigger model needs bigger data\" — classic ML principle",
    ]),
    ("2.  No data augmentation in this study", [
        "Augmentation (flip, rotation, intensity) disproportionately benefits over-parameterised models",
        "Follow-up: add augmentation + scale to 150–369 patients",
    ]),
    ("3.  Same hyperparameters for all models", [
        "Transformers typically need lower LR + warmup (e.g., 5e-4 with linear warmup)",
        "Uniform 1e-3 may have hurt Hybrid convergence",
    ]),
]
y = 1.65
for title, lines in boxes:
    add_textbox(disc, Inches(0.5), Inches(y), Inches(9.0), Inches(0.4),
                [{"text": title, "size": 14, "bold": True, "color": ACCENT}])
    add_textbox(disc, Inches(0.7), Inches(y + 0.4), Inches(8.8), Inches(1.1),
                [{"text": "• " + l, "size": 11} for l in lines])
    y += 1.55
add_textbox(disc, Inches(0.5), Inches(6.4), Inches(9.0), Inches(0.7),
            [{"text": "Takeaway:  Architectural progress is real, but it requires the right experimental conditions to manifest. Honesty in reporting > forced-positive results.",
              "size": 11, "italic": True, "bold": True, "color": NAVY}])
move_slide(prs, len(prs.slides) - 1, 14)


# ===== Slide index 15: Conclusion & Future Work =====
s = prs.slides[15]
cp = find_content_placeholder(s)
if cp is not None:
    remove_shape(cp)
# also remove decorative inner Rectangle 3 ("Conclusion" label) so we have free space
inner = find_shape(s, "Rectangle 3")
if inner is not None:
    remove_shape(inner)
add_textbox(s, Inches(0.5), Inches(1.05), Inches(9.0), Inches(0.5),
            [{"text": "Contributions & Roadmap", "size": 18, "bold": True, "color": NAVY}])
# Achieved
add_textbox(s, Inches(0.5), Inches(1.55), Inches(4.5), Inches(2.5), [
    {"text": "Achieved", "size": 14, "bold": True, "color": ACCENT},
    {"text": "• 3 progressively complex 2D U-Net variants on BraTS 2020", "size": 11},
    {"text": "• Best: Dice 0.8370 / IoU 0.7715 (UNet2D + DiceBCE)", "size": 11},
    {"text": "• Empirical comparison of DiceBCE vs Focal Tversky", "size": 11},
    {"text": "• Honest characterisation of capacity-vs-data tradeoff", "size": 11},
])
# Why matters
add_textbox(s, Inches(0.5), Inches(4.0), Inches(4.5), Inches(2.5), [
    {"text": "Why this matters", "size": 14, "bold": True, "color": ACCENT},
    {"text": "• 15 min on T4 vs days on A100 for nnU-Net", "size": 11},
    {"text": "• Open-source, reproducible end-to-end (Colab + GitHub)", "size": 11},
    {"text": "• Clean ablation framework for adding components", "size": 11},
])
# Future work table
add_textbox(s, Inches(5.2), Inches(1.55), Inches(4.6), Inches(0.4),
            [{"text": "Future work", "size": 14, "bold": True, "color": ACCENT}])
fw = [
    ["Improvement", "Δ Dice"],
    ["Scale to 150–369 patients + augmentation", "+2–4%"],
    ["3D U-Net (volumetric context)", "+5–7%"],
    ["Test-time augmentation (TTA)", "+0.5–1%"],
    ["3-class output (NCR/ED/ET) vs binary", "richer"],
    ["Lower LR + warmup for transformer", "+1–2%"],
]
add_table(s, Inches(5.2), Inches(2.0), Inches(4.6), Inches(3.5), fw,
          col_widths=[Inches(3.2), Inches(1.4)],
          body_size=10, header_size=11, header_fill=ACCENT)


# ===== Slide index 16: References =====
s = prs.slides[16]
cp = find_content_placeholder(s)
refs = [
    "[1]  Ronneberger, O., Fischer, P., & Brox, T. (2015). U-Net: Convolutional Networks for Biomedical Image Segmentation. MICCAI.",
    "[2]  Milletari, F., Navab, N., & Ahmadi, S.-A. (2016). V-Net: Fully Convolutional Neural Networks for Volumetric Medical Image Segmentation. 3DV.",
    "[3]  Oktay, O. et al. (2018). Attention U-Net: Learning Where to Look for the Pancreas. MIDL.",
    "[4]  Wang, W. et al. (2021). TransBTS: Multimodal Brain Tumor Segmentation Using Transformer. MICCAI.",
    "[5]  Hatamizadeh, A. et al. (2022). Swin UNETR. Brainlesion Workshop, MICCAI.",
    "[6]  Isensee, F. et al. (2021). nnU-Net: A Self-Configuring Method for Deep Learning-Based Biomedical Image Segmentation. Nature Methods.",
    "[7]  Abraham, N., & Khan, N. M. (2019). A Novel Focal Tversky Loss Function with Improved Attention U-Net for Lesion Segmentation. ISBI.",
    "[8]  Salehi, S. S. M., Erdogmus, D., & Gholipour, A. (2017). Tversky Loss Function for Image Segmentation. MICCAI MLMI.",
    "[9]  Vaswani, A. et al. (2017). Attention is All You Need. NeurIPS.",
    "[10] Menze, B. H. et al. (2015). The Multimodal Brain Tumor Image Segmentation Benchmark (BraTS). IEEE TMI.",
]
set_paragraphs(cp.text_frame,
               [{"text": r, "size": 11} for r in refs])


# ===== Slide 17 (last): Questions slide is unchanged =====


prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH}")
print(f"  Size: {OUT_PATH.stat().st_size/1024:.1f} KB")
print(f"  Slides: {len(prs.slides)}")
